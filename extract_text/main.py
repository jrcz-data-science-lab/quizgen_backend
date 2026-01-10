from fastapi import FastAPI, UploadFile, File
import pypdf
from docx import Document
from pptx import Presentation
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI()

def extract_text_from_pdf(path):
    text = ""
    reader = pypdf.PdfReader(path)
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_pptx(path):
    presentation = Presentation(path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    filename = f"temp.{file.filename.split('.')[-1]}"
    contents = await file.read()

    # Save uploaded file
    with open(filename, "wb") as f:
        f.write(contents)

    ext = filename.split(".")[-1]

    # Extract based on file type
    if ext == "pdf":
        extracted = extract_text_from_pdf(filename)
    elif ext == "docx":
        extracted = extract_text_from_docx(filename)
    elif ext == "pptx":
        extracted = extract_text_from_pptx(filename)
    else:
        return {"error": "Unsupported file type"}

    return {
        "extracted_text": extracted
    }
